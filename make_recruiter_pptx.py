#!/usr/bin/env python3
"""Generate a 12-slide PowerPoint overview for the async FIFO project.

Dependency:
  python-pptx

Example:
  python3 make_recruiter_pptx.py
"""

from __future__ import annotations

from pathlib import Path

try:
    from PIL import Image
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - exercised manually.
    raise SystemExit(
        "Missing dependency. Install with: python3 -m pip install python-pptx\n"
        f"Import error: {exc}"
    ) from exc


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "AsyncFIFO_CDC_Verification_Overview.pptx"

ASSETS = {
    "flag_overview": ROOT / "Post_M5/UVM/flag_debug_waveforms.png",
    "write_threshold": ROOT / "Post_M5/UVM/flag_debug_write_thresholds.png",
    "read_threshold": ROOT / "Post_M5/UVM/flag_debug_read_thresholds.png",
    "coverage": ROOT / "Post_M5/docs/coverage_summary.png",
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(12, 32, 52)
BLUE = RGBColor(37, 99, 166)
LIGHT_BLUE = RGBColor(232, 241, 250)
GREEN = RGBColor(42, 125, 74)
LIGHT_GREEN = RGBColor(231, 245, 236)
GOLD = RGBColor(190, 137, 34)
LIGHT_GOLD = RGBColor(252, 245, 224)
GRAY = RGBColor(88, 96, 105)
DARK = RGBColor(31, 35, 40)
LIGHT_GRAY = RGBColor(245, 247, 250)
MID_GRAY = RGBColor(218, 224, 231)
WHITE = RGBColor(255, 255, 255)


def validate_assets() -> None:
    missing = [str(path.relative_to(ROOT)) for path in ASSETS.values() if not path.exists()]
    if missing:
        raise SystemExit("Missing required image assets:\n" + "\n".join(missing))


def set_run_font(run, size: int, color: RGBColor = DARK, bold: bool = False) -> None:
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    size: int = 18,
    color: RGBColor = DARK,
    bold: bool = False,
    align_center: bool = False,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    if align_center:
        p.alignment = 1
    for run in p.runs:
        set_run_font(run, size, color, bold)
    return box


def add_bullets(slide, items: list[str], left, top, width, height, size: int = 16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, Inches(0.55), Inches(0.32), Inches(12.25), Inches(0.45), 26, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.58), Inches(0.82), Inches(12.0), Inches(0.35), 12, GRAY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.18), Inches(12.25), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_footer(slide, index: int) -> None:
    add_text(
        slide,
        f"Asynchronous FIFO CDC Design and Verification Study | {index:02d}",
        Inches(0.55),
        Inches(7.08),
        Inches(12.25),
        Inches(0.25),
        9,
        GRAY,
    )


def add_card(slide, title: str, body: str, left, top, width, height, fill=LIGHT_BLUE, accent=BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title
    for run in p.runs:
        set_run_font(run, 15, NAVY, True)
    p = tf.add_paragraph()
    p.text = body
    p.space_before = Pt(6)
    for run in p.runs:
        set_run_font(run, 12, DARK)


def add_metric(slide, value: str, label: str, left, top, width, height, color=BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = 1
    for run in p.runs:
        set_run_font(run, 21, color, True)
    p = tf.add_paragraph()
    p.text = label
    p.alignment = 1
    for run in p.runs:
        set_run_font(run, 10, GRAY)


def add_fitted_image(slide, image_path: Path, left, top, width, height) -> None:
    with Image.open(image_path) as img:
        img_w, img_h = img.size
    img_ratio = img_w / img_h
    box_ratio = width / height
    if img_ratio >= box_ratio:
        pic_w = width
        pic_h = width / img_ratio
    else:
        pic_h = height
        pic_w = height * img_ratio
    pic_left = left + (width - pic_w) / 2
    pic_top = top + (height - pic_h) / 2
    slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_w, height=pic_h)


def add_table(slide, rows: list[list[str]], left, top, width, height) -> None:
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            fill = NAVY if r == 0 else (LIGHT_GRAY if r % 2 == 0 else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    set_run_font(run, 10 if r else 11, WHITE if r == 0 else DARK, r == 0)


def add_block(slide, text: str, left, top, width, height, fill=WHITE, line=BLUE, size=12) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = 1
    for run in p.runs:
        set_run_font(run, size, NAVY, True)


def add_arrow(slide, left, top, width, height, fill=BLUE) -> None:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = fill
    arrow.line.fill.background()


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation, index: int, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 252, 255)
    bg.line.fill.background()
    add_title(slide, title, subtitle)
    add_footer(slide, index)
    return slide


def build_deck() -> Presentation:
    prs = new_deck()

    # 1. Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 251, 255)
    bg.line.fill.background()
    add_text(slide, "Asynchronous FIFO CDC Design", Inches(0.7), Inches(0.72), Inches(12), Inches(0.55), 34, NAVY, True)
    add_text(slide, "SystemVerilog/UVM verification study with farm-backed evidence", Inches(0.72), Inches(1.32), Inches(12), Inches(0.38), 16, GRAY)
    add_fitted_image(slide, ASSETS["flag_overview"], Inches(0.65), Inches(2.0), Inches(12.1), Inches(3.5))
    add_metric(slide, "Post_M5/UVM", "canonical target", Inches(0.8), Inches(6.0), Inches(2.5), Inches(0.72), BLUE)
    add_metric(slide, "0", "UVM errors/fatals", Inches(3.65), Inches(6.0), Inches(2.2), Inches(0.72), GREEN)
    add_metric(slide, "100%", "coverage goals hit", Inches(6.25), Inches(6.0), Inches(2.2), Inches(0.72), GREEN)
    add_metric(slide, "Questa", "farm-verified", Inches(8.85), Inches(6.0), Inches(2.2), Inches(0.72), GOLD)
    add_footer(slide, 1)

    # 2. Problem statement
    slide = blank_slide(prs, 2, "Problem Statement", "Safe data transfer across unrelated write and read clocks")
    add_block(slide, "Write domain\n80 MHz\n12.5 ns", Inches(0.8), Inches(2.0), Inches(2.3), Inches(1.25), LIGHT_BLUE)
    add_arrow(slide, Inches(3.25), Inches(2.42), Inches(1.25), Inches(0.35))
    add_block(slide, "Async FIFO\n64 entries\n8-bit data", Inches(4.75), Inches(1.85), Inches(2.6), Inches(1.55), LIGHT_GOLD, GOLD)
    add_arrow(slide, Inches(7.62), Inches(2.42), Inches(1.25), Inches(0.35))
    add_block(slide, "Read domain\n50 MHz\n20 ns", Inches(9.1), Inches(2.0), Inches(2.3), Inches(1.25), LIGHT_GREEN, GREEN)
    add_bullets(
        slide,
        [
            "Preserve FIFO ordering while the two clocks advance independently.",
            "Prevent overflow writes and underflow reads without sharing unsafe dual-clock state.",
            "Generate full, empty, half-full, and half-empty flags with threshold-accurate timing.",
            "Capture proof in simulator transcripts and reviewable waveform artifacts.",
        ],
        Inches(1.05),
        Inches(4.1),
        Inches(11.2),
        Inches(1.7),
        15,
    )

    # 3. Design objectives
    slide = blank_slide(prs, 3, "Design Objectives", "Implementation and verification goals")
    add_card(slide, "CDC correctness", "Synchronize pointers safely across unrelated domains.", Inches(0.75), Inches(1.55), Inches(2.85), Inches(1.45), LIGHT_BLUE)
    add_card(slide, "Data integrity", "Maintain ordering through wraparound, backpressure, and resets.", Inches(3.9), Inches(1.55), Inches(2.85), Inches(1.45), LIGHT_GREEN, GREEN)
    add_card(slide, "Flag accuracy", "Assert full, empty, half-full, and half-empty at the intended edges.", Inches(7.05), Inches(1.55), Inches(2.85), Inches(1.45), LIGHT_GOLD, GOLD)
    add_card(slide, "Evidence quality", "Commit farm transcripts, coverage reports, and readable waveform images.", Inches(10.2), Inches(1.55), Inches(2.35), Inches(1.45), RGBColor(241, 238, 249), RGBColor(111, 66, 193))
    add_text(slide, "Primary implementation principle", Inches(0.95), Inches(3.75), Inches(5.5), Inches(0.32), 18, NAVY, True)
    add_bullets(
        slide,
        [
            "Use local binary pointers for arithmetic.",
            "Send only Gray-coded pointers across clock boundaries.",
            "Place both synchronizer flops in the destination domain.",
            "Compute occupancy locally from the next pointer and synchronized remote pointer.",
        ],
        Inches(0.95),
        Inches(4.25),
        Inches(5.6),
        Inches(1.7),
        14,
    )
    add_text(slide, "Verification principle", Inches(7.0), Inches(3.75), Inches(5.4), Inches(0.32), 18, NAVY, True)
    add_bullets(
        slide,
        [
            "Drive before the DUT sampling edge.",
            "Monitor accepted transfers, not raw attempts.",
            "Check data with an independent queue reference model.",
            "Use directed tests for exact flag thresholds.",
        ],
        Inches(7.0),
        Inches(4.25),
        Inches(5.5),
        Inches(1.7),
        14,
    )

    # 4. DUT architecture
    slide = blank_slide(prs, 4, "DUT Architecture", "Five-block async FIFO structure")
    add_block(slide, "write_pointer\nbinary -> Gray\nwFull, wHalfFull", Inches(0.8), Inches(1.55), Inches(2.55), Inches(1.05), LIGHT_BLUE)
    add_block(slide, "sync_w2r\n2 flops in rclk", Inches(3.75), Inches(1.55), Inches(2.15), Inches(1.05), WHITE)
    add_block(slide, "read_pointer\nbinary -> Gray\nrEmpty, rHalfEmpty", Inches(6.3), Inches(1.55), Inches(2.55), Inches(1.05), LIGHT_GREEN, GREEN)
    add_block(slide, "fifo_memory\n64 x 8 storage", Inches(4.65), Inches(3.15), Inches(3.0), Inches(1.05), LIGHT_GOLD, GOLD)
    add_block(slide, "sync_r2w\n2 flops in wclk", Inches(3.75), Inches(4.75), Inches(2.15), Inches(1.05), WHITE)
    add_arrow(slide, Inches(3.35), Inches(1.88), Inches(0.35), Inches(0.22))
    add_arrow(slide, Inches(5.92), Inches(1.88), Inches(0.35), Inches(0.22))
    add_arrow(slide, Inches(6.6), Inches(4.45), Inches(0.35), Inches(0.22), GREEN)
    add_arrow(slide, Inches(3.35), Inches(5.08), Inches(0.35), Inches(0.22), GREEN)
    add_bullets(
        slide,
        [
            "Memory access is gated by local full/empty protection.",
            "Pointer width is one bit wider than the RAM address.",
            "Gray-coded pointers are the only cross-domain state.",
            "Each flag is registered in the domain that consumes it.",
        ],
        Inches(9.35),
        Inches(1.6),
        Inches(3.0),
        Inches(4.1),
        14,
    )

    # 5. CDC decisions
    slide = blank_slide(prs, 5, "CDC Design Decisions", "Why the final RTL avoids common async FIFO traps")
    add_card(slide, "Binary locally", "Increment, address, and occupancy math remain simple within one clock domain.", Inches(0.75), Inches(1.55), Inches(3.65), Inches(1.32), LIGHT_BLUE)
    add_card(slide, "Gray across domains", "Only one pointer bit changes per increment before synchronization.", Inches(4.85), Inches(1.55), Inches(3.65), Inches(1.32), LIGHT_GREEN, GREEN)
    add_card(slide, "Extra MSB", "Same address can mean empty or full depending on pointer wrap state.", Inches(8.95), Inches(1.55), Inches(3.3), Inches(1.32), LIGHT_GOLD, GOLD)
    add_text(slide, "Rejected pattern", Inches(0.95), Inches(3.55), Inches(3.7), Inches(0.35), 18, NAVY, True)
    add_block(slide, "single fifo_count\nwritten by wclk and rclk", Inches(0.95), Inches(4.15), Inches(3.4), Inches(1.15), RGBColor(253, 237, 237), RGBColor(180, 35, 24))
    add_text(slide, "Final pattern", Inches(5.0), Inches(3.55), Inches(3.7), Inches(0.35), 18, NAVY, True)
    add_block(slide, "next local pointer\nminus synced remote pointer", Inches(5.0), Inches(4.15), Inches(3.6), Inches(1.15), LIGHT_BLUE, BLUE)
    add_text(slide, "Result", Inches(9.2), Inches(3.55), Inches(3.1), Inches(0.35), 18, NAVY, True)
    add_block(slide, "CDC-safe occupancy\nper consuming domain", Inches(9.2), Inches(4.15), Inches(3.05), Inches(1.15), LIGHT_GREEN, GREEN)

    # 6. Flag timing innovation
    slide = blank_slide(prs, 6, "Flag Timing Innovation", "Write-side half-full and full proof")
    add_fitted_image(slide, ASSETS["write_threshold"], Inches(0.55), Inches(1.42), Inches(8.45), Inches(4.6))
    add_bullets(
        slide,
        [
            "Half-full uses `binary_wptr_next`, not the prior registered pointer.",
            "`wHalfFull` asserts when the accepted write reaches 32 entries.",
            "`wFull` asserts when the accepted write reaches 64 entries.",
            "The waveform is generated from the directed farm VCD.",
        ],
        Inches(9.25),
        Inches(1.65),
        Inches(3.3),
        Inches(3.3),
        14,
    )
    add_metric(slide, "32", "half-full threshold", Inches(9.35), Inches(5.15), Inches(1.35), Inches(0.68), BLUE)
    add_metric(slide, "64", "full threshold", Inches(10.95), Inches(5.15), Inches(1.35), Inches(0.68), GREEN)

    # 7. Read-side evidence
    slide = blank_slide(prs, 7, "Read-Side Flag Evidence", "Half-empty and empty threshold proof")
    add_fitted_image(slide, ASSETS["read_threshold"], Inches(0.55), Inches(1.42), Inches(8.45), Inches(4.6))
    add_bullets(
        slide,
        [
            "Half-empty uses `binary_rptr_next` for same-edge threshold behavior.",
            "`rHalfEmpty` asserts when occupancy drops to 32 entries.",
            "`rEmpty` asserts when the final accepted read drains the FIFO.",
            "Overflow and underflow attempts are tested without corrupting the scoreboard model.",
        ],
        Inches(9.25),
        Inches(1.65),
        Inches(3.3),
        Inches(3.3),
        14,
    )
    add_metric(slide, "32", "half-empty threshold", Inches(9.35), Inches(5.15), Inches(1.35), Inches(0.68), GOLD)
    add_metric(slide, "0", "empty occupancy", Inches(10.95), Inches(5.15), Inches(1.35), Inches(0.68), GREEN)

    # 8. Verification architecture
    slide = blank_slide(prs, 8, "UVM Verification Architecture", "Independent agents for independent clock domains")
    add_block(slide, "fifo_random_test\nor fifo_base_test", Inches(0.75), Inches(1.55), Inches(2.4), Inches(0.8), LIGHT_GOLD, GOLD)
    add_block(slide, "fifo_env", Inches(3.55), Inches(1.55), Inches(1.8), Inches(0.8), WHITE)
    add_block(slide, "write_agent\nsequencer + driver + monitor", Inches(1.05), Inches(3.05), Inches(3.05), Inches(1.15), LIGHT_BLUE)
    add_block(slide, "read_agent\nsequencer + driver + monitor", Inches(4.55), Inches(3.05), Inches(3.05), Inches(1.15), LIGHT_GREEN, GREEN)
    add_block(slide, "fifo_scoreboard\nqueue reference model", Inches(8.25), Inches(3.05), Inches(3.2), Inches(1.15), LIGHT_GOLD, GOLD)
    add_arrow(slide, Inches(3.16), Inches(1.83), Inches(0.35), Inches(0.22))
    add_arrow(slide, Inches(4.12), Inches(4.55), Inches(0.35), Inches(0.22))
    add_arrow(slide, Inches(7.68), Inches(4.55), Inches(0.35), Inches(0.22), GREEN)
    add_bullets(
        slide,
        [
            "Separate read and write paths reflect the real CDC split.",
            "Drivers produce stable stimulus before the active edge.",
            "Monitors publish only accepted transfers.",
            "The scoreboard checks ordering independently of DUT internals.",
        ],
        Inches(0.95),
        Inches(5.15),
        Inches(11.2),
        Inches(1.0),
        14,
    )

    # 9. Scoreboard and monitors
    slide = blank_slide(prs, 9, "Scoreboard And Monitor Strategy", "Accepted-transfer modeling with reset-aware checking")
    add_block(slide, "Write monitor\nwinc && !wFull", Inches(0.85), Inches(1.65), Inches(2.35), Inches(0.9), LIGHT_BLUE)
    add_arrow(slide, Inches(3.35), Inches(1.94), Inches(0.55), Inches(0.25))
    add_block(slide, "Expected queue\npush wData", Inches(4.1), Inches(1.65), Inches(2.25), Inches(0.9), WHITE)
    add_block(slide, "Read monitor\nrinc && !rEmpty", Inches(0.85), Inches(3.25), Inches(2.35), Inches(0.9), LIGHT_GREEN, GREEN)
    add_arrow(slide, Inches(3.35), Inches(3.54), Inches(0.55), Inches(0.25), GREEN)
    add_block(slide, "Compare rData\npop expected", Inches(4.1), Inches(3.25), Inches(2.25), Inches(0.9), WHITE, GREEN)
    add_block(slide, "Reset monitor\nflush queue", Inches(0.85), Inches(4.85), Inches(2.35), Inches(0.9), LIGHT_GOLD, GOLD)
    add_arrow(slide, Inches(3.35), Inches(5.14), Inches(0.55), Inches(0.25), GOLD)
    add_block(slide, "Aligned model\nafter reset", Inches(4.1), Inches(4.85), Inches(2.25), Inches(0.9), WHITE, GOLD)
    add_text(slide, "Farm scoreboard summary", Inches(7.2), Inches(1.55), Inches(4.8), Inches(0.35), 18, NAVY, True)
    add_table(
        slide,
        [
            ["Metric", "Result"],
            ["Writes observed", "659"],
            ["Reads observed", "596"],
            ["Reset flushes", "2"],
            ["Residual queue", "0"],
            ["Mismatches / errors", "0"],
            ["Verdict", "PASSED"],
        ],
        Inches(7.2),
        Inches(2.05),
        Inches(4.75),
        Inches(3.6),
    )

    # 10. Coverage plan
    slide = blank_slide(prs, 10, "Coverage Plan", "DUT-scoped code coverage plus functional covergroups")
    add_fitted_image(slide, ASSETS["coverage"], Inches(0.55), Inches(1.42), Inches(7.6), Inches(4.45))
    add_bullets(
        slide,
        [
            "DUT-focused coverage avoids inflating results with testbench code.",
            "Code coverage: statements, branches, expressions, conditions.",
            "Functional coverage: flags, data patterns, reset, bursts, idle cycles, throughput.",
            "Final result: 100% DUT-filtered code coverage and 100% covergroup coverage.",
        ],
        Inches(8.55),
        Inches(1.65),
        Inches(3.75),
        Inches(3.0),
        14,
    )
    add_metric(slide, "100%", "DUT coverage", Inches(8.7), Inches(5.15), Inches(1.5), Inches(0.68), GREEN)
    add_metric(slide, "9", "covergroups", Inches(10.45), Inches(5.15), Inches(1.5), Inches(0.68), BLUE)

    # 11. Bug injection
    slide = blank_slide(prs, 11, "Bug Injection And Robustness", "Negative testing confirms the checker can fail for real defects")
    add_table(
        slide,
        [
            ["Macro", "Injected defect", "Verification purpose"],
            ["WDATA_CORRUPTION_BUG", "Corrupt write data", "Scoreboard catches data mismatch"],
            ["SYNC_BUG", "Remove second sync flop", "Exercise CDC assumptions"],
            ["RPTR_BUG", "Use binary read pointer", "Test Gray encoding checks"],
            ["WPTR_FULLFLAG_BUG", "Naive full equality", "Validate wrap-aware full logic"],
        ],
        Inches(0.75),
        Inches(1.55),
        Inches(11.8),
        Inches(2.3),
    )
    add_bullets(
        slide,
        [
            "Default build keeps all defect macros disabled.",
            "Each macro is intended to be enabled independently.",
            "The goal is not only to pass the nominal design, but to prove that meaningful bugs are detectable.",
        ],
        Inches(1.0),
        Inches(4.35),
        Inches(11.0),
        Inches(1.4),
        15,
    )

    # 12. Final results
    slide = blank_slide(prs, 12, "Final Results And Review Path", "Concise evidence for project evaluation")
    add_metric(slide, "0", "UVM warnings/errors/fatals", Inches(0.75), Inches(1.55), Inches(2.55), Inches(0.9), GREEN)
    add_metric(slide, "100%", "DUT-focused coverage", Inches(3.65), Inches(1.55), Inches(2.55), Inches(0.9), GREEN)
    add_metric(slide, "100%", "covergroup coverage", Inches(6.55), Inches(1.55), Inches(2.55), Inches(0.9), GREEN)
    add_metric(slide, "PASS", "directed flag thresholds", Inches(9.45), Inches(1.55), Inches(2.55), Inches(0.9), BLUE)
    add_text(slide, "Recommended review path", Inches(0.9), Inches(3.2), Inches(5.0), Inches(0.35), 18, NAVY, True)
    add_bullets(
        slide,
        [
            "Post_M5/UVM/async_fifo.sv",
            "Post_M5/UVM/async_fifo_scoreboard.sv",
            "Post_M5/UVM/async_fifo_coverage.sv",
            "Post_M5/UVM/MANIFEST.txt",
        ],
        Inches(0.9),
        Inches(3.75),
        Inches(5.1),
        Inches(1.5),
        14,
    )
    add_text(slide, "Closing summary", Inches(6.55), Inches(3.2), Inches(5.6), Inches(0.35), 18, NAVY, True)
    add_bullets(
        slide,
        [
            "CDC-safe FIFO RTL with Gray-coded pointer synchronization.",
            "UVM environment with independent agents and reset-aware checking.",
            "Directed waveform evidence for threshold-accurate flags.",
            "Farm-backed transcripts and coverage artifacts committed with the repo.",
        ],
        Inches(6.55),
        Inches(3.75),
        Inches(5.8),
        Inches(1.6),
        14,
    )

    return prs


def main() -> None:
    validate_assets()
    prs = build_deck()
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(ROOT)} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
